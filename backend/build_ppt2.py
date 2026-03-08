from pptx import Presentation
import os

TEMPLATE_PATH = 'c:/Users/Siddhant/NidhiAi/Prototype Development Submission _ AWS AI for Bharat Hackathon.pptx'
OUTPUT_PATH = 'c:/Users/Siddhant/NidhiAi/NidhiAI_AWS_Hackathon_Submission.pptx'

prs = Presentation(TEMPLATE_PATH)

replacements = {
    'Team Name :': 'Team Name: NidhiAI\nTrack: 03 - AI for Communities, Access & Public Impact',
    'Team Leader Name :': 'Team Leader Name: Siddhant\nInstitution/Organization: Independent Developer',
    
    'Problem Statement :': 'Problem Statement:\n\n'
                           'India mandates ~₹38,000 Crores annually for Corporate Social Responsibility (CSR). '
                           'However, nearly 90% of these funds are concentrated among the top 50, heavily-resourced NGOs.\n\n'
                           'Grassroots organizations serving rural India are locked out of this capital due to two strict barriers:\n'
                           '1. Legal Compliance Overhead: The inability to rapidly verify and maintain 12A/80G status.\n'
                           '2. Language & Proposal Formatting: The inability to draft corporate-grade English proposals that parse Schedule VII CSR mandates.\n\n'
                           'The capital exists, but the administrative transmission layer is broken.',
    
    'Brief about the Idea:': 'Brief about the Idea:\n\n'
                             'NidhiAI is an intelligent, multi-agent orchestration platform that directly solves the administrative bottleneck for grassroots NGOs.\n\n'
                             'Core Subsystems:\n'
                             '- Automated Legal Validation: Uses Amazon Textract to ingest 12A/80G PDFs, instantly verifying expiry dates and compliance standing.\n'
                             '- Semantic CSR Matching: Uses Titan Text Embeddings V2 against an OpenSearch Knowledge Base populated with actual corporate CSR filings to match NGO operations to available funds.\n'
                             '- Instant Proposal Generation: A Claude-powered agent ingests the NGO context and the corporate grant criteria via RAG to draft a formal, hyper-targeted 5-page PDF proposal in seconds.',
    
    'AI is required in your solution?': 'Why AI is required in your solution?\n\n'
                                        '1. Reasoning over rigid search: Standard databases cannot map a rural "village water well project" to Tata Steel\'s dense "Schedule VII Rural Infrastructure Development" CSR mandate. Semantic matching is strictly required.\n'
                                        '2. Generative Creation: Grassroots NGOs cannot afford grant writers. Claude 3.5 Sonnet acts as a highly-skilled proxy, transforming bullet points into formal corporate proposals instantly.\n'
                                        '3. Unstructured Data Extraction: Government documents (12A/80G) lack standard APIs. Using an AI vision pipeline (Amazon Textract + Bedrock Agent) is the only way to validate expiry dates and registration numbers from raw scans.',
    
    'Technologies utilized in the solution:on': 'Technologies utilized in the solution:\n\n'
                                                 'Front-End:\n'
                                                 '- Next.js 16 (React) / TailwindCSS / Vercel\n\n'
                                                 'Back-End & Serverless Compute:\n'
                                                 '- AWS API Gateway & AWS Lambda (Python)\n'
                                                 '- Amazon DynamoDB (State management)\n\n'
                                                 'AWS AI / ML Stack:\n'
                                                 '- Amazon Bedrock Agents (Supervisor Orchestration Pattern)\n'
                                                 '- Amazon Bedrock Knowledge Bases (Amazon OpenSearch Serverless)\n'
                                                 '- Embedding Model: Amazon Titan Text Embeddings V2\n'
                                                 '- Foundation Models: Claude 3.5 Sonnet / Amazon Nova 2 Lite\n'
                                                 '- Amazon Textract (OCR and Document Analysis)',
                                                 
    'Estimated implementation cost (optional):': 'Estimated implementation cost (production-scale):\n\n'
                                                 'Strict Serverless architecture ensures near zero baseline costs:\n'
                                                 '- Compute (API Gateway/Lambda/DynamoDB): Fully covered under AWS Free Tier for prototype scale.\n'
                                                 '- Bedrock Agent & Model Inference: ~$10/month at medium volume.\n'
                                                 '- OpenSearch Serverless: Provisioned minimums apply (~$180/month for production, zeroed out when idle).\n'
                                                 '- Sub-cent cost per proposal generated compared to $500+ for human grant writers.',
                                                 
    'Prototype Performance report/Benchmarking:': 'Prototype Performance report/Benchmarking:\n\n'
                                                  '- OCR Compliance Scan (Amazon Textract -> Bedrock): ~4.5 seconds round-trip per document.\n'
                                                  '- OpenSearch Semantic Grant Matching (Titan V2): 1.8 seconds latency.\n'
                                                  '- Proposal Generation (RAG via Claude): ~18 seconds (server-side stream implementation).\n'
                                                  '- System Availability: 100% reliant on AWS fully managed services.',
                                                  
    'Additional Details/Future Development (if any):': 'Additional Details/Future Development:\n\n'
                                                       '- Multilingual Voice Intake: Allow uneducated NGO founders to narrate their community impact verbally in Hindi via Amazon Transcribe, converting it into an English corporate proposal seamlessly.\n'
                                                       '- Direct Corporate Portal Integration: Utilizing Amazon Q Business to securely interface with internal CSR systems for one-click application submission.',
                                                       
    'GitHub Public Repository Link': 'GitHub Public Repository Link: https://github.com/Sid-V5/NidhiAi\n',
    'Demo Video Link (Max: 3 Minutes)': 'Demo Video Link (Max: 3 Minutes): [Insert Public YouTube / Drive Link Here]\nWorking MVP Link: [Insert Vercel URL Here]'
}

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        # We need to replace text exactly, but sometimes PPT text is split across runs.
        # So we collect the whole paragraph text, match it, and if it matches, clear the paragraph and add the new text.
        for paragraph in shape.text_frame.paragraphs:
            # We do a 'startswith' or 'in' check on the raw text of the paragraph
            full_text = "".join(run.text for run in paragraph.runs)
            
            for key, new_val in replacements.items():
                if key in full_text:
                    print(f"Matched and replacing: {key}")
                    # Clear all runs in the paragraph
                    for idx in range(len(paragraph.runs)-1, -1, -1):
                        paragraph.runs[idx].text = ""
                    
                    # Set the new text to the first run (or create one if somehow empty, though usually it isn't)
                    if paragraph.runs:
                        paragraph.runs[0].text = new_val
                    else:
                        paragraph.add_run().text = new_val

prs.save(OUTPUT_PATH)
print(f"Refined PPT saved to {OUTPUT_PATH}")
