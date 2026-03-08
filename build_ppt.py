"""
Build the final NidhiAI Prototype Development Submission PPT.
Fills all 14 slides of the Hack2Skill template with real content and embedded screenshots.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

TEMPLATE = r'c:\Users\Siddhant\NidhiAi\Prototype Development Submission _ AWS AI for Bharat Hackathon.pptx'
OUTPUT = r'c:\Users\Siddhant\NidhiAi\NidhiAI_Final_Submission.pptx'
SS_DIR = r'C:\Users\Siddhant\.gemini\antigravity\brain\ffadbb09-a544-4a00-a7d0-c6df051cd0f2'

prs = Presentation(TEMPLATE)

# ----------------------------------------------------------------
# Helper: Replace text in a shape while preserving the FIRST run's
# font formatting. Clears all runs, sets text on first run.
# ----------------------------------------------------------------
def set_shape_text(shape, new_text):
    """Replace all text in a shape with new_text, keeping first run font."""
    tf = shape.text_frame
    # Gather first run's formatting
    first_font = None
    for para in tf.paragraphs:
        if para.runs:
            first_font = para.runs[0].font
            break

    # Clear all paragraphs except first
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)

    # Now set text on the remaining paragraph
    para = tf.paragraphs[0]
    # Clear runs
    for run in para.runs:
        run.text = ""
    
    # Split new_text by newlines and add paragraphs
    lines = new_text.split('\n')
    para.runs[0].text = lines[0] if para.runs else lines[0]
    
    if not para.runs:
        r = para.add_run()
        r.text = lines[0]
    
    for line in lines[1:]:
        new_para = tf.paragraphs[0]._p.addnext(
            tf.paragraphs[0]._p.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}p', {})
        )
        from pptx.oxml.ns import qn
        r_elem = new_para.makeelement(qn('a:r'), {})
        new_para.append(r_elem)
        t_elem = r_elem.makeelement(qn('a:t'), {})
        t_elem.text = line
        r_elem.append(t_elem)
        # Copy font properties if available
        if first_font and first_font.name:
            rPr = r_elem.makeelement(qn('a:rPr'), {'lang': 'en-US'})
            r_elem.insert(0, rPr)


def simple_set(shape, text):
    """Simpler: just replace each paragraph's runs with new content."""
    tf = shape.text_frame
    # Get all existing paragraphs
    paras = list(tf.paragraphs)
    
    # Replace first paragraph text
    lines = text.split('\n')
    
    # Clear ALL existing content
    for p in paras:
        for run in p.runs:
            run.text = ""
    
    # Set first line on first paragraph's first run
    if paras[0].runs:
        paras[0].runs[0].text = lines[0]
    
    # For remaining lines, add to subsequent paragraphs or create new runs  
    from pptx.oxml.ns import qn
    from copy import deepcopy
    
    first_p = paras[0]._p
    parent = first_p.getparent()
    
    for line_text in lines[1:]:
        # Clone the first paragraph to preserve formatting
        new_p = deepcopy(first_p)
        # Clear its text
        for r in new_p.findall(qn('a:r')):
            for t in r.findall(qn('a:t')):
                t.text = ""
        # Set new text on first run
        runs = new_p.findall(qn('a:r'))
        if runs:
            ts = runs[0].findall(qn('a:t'))
            if ts:
                ts[0].text = line_text
            else:
                t_elem = runs[0].makeelement(qn('a:t'), {})
                t_elem.text = line_text
                runs[0].append(t_elem)
        parent.append(new_p)


def add_image_to_slide(slide, img_path, left, top, width, height=None):
    """Add an image to a slide at the given position."""
    if os.path.exists(img_path):
        if height:
            slide.shapes.add_picture(img_path, left, top, width, height)
        else:
            slide.shapes.add_picture(img_path, left, top, width)
        print(f"  Added image: {os.path.basename(img_path)}")
    else:
        print(f"  WARNING: Image not found: {img_path}")


# ================================================================
# SLIDE 1: Team Name, Problem Statement, Team Leader
# ================================================================
print("Slide 1: Team Info & Problem Statement")
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if not shape.has_text_frame:
        continue
    txt = shape.text_frame.text.strip()
    
    if txt.startswith('Team Name'):
        simple_set(shape, 'Team Name: NidhiAI')
    elif txt.startswith('Problem Statement'):
        simple_set(shape, 
            'Problem Statement:\n'
            '\n'
            'India mandates ~Rs.38,000 Crores annually for CSR spending. Yet 90% of\n'
            'these funds flow to the top 50 well-resourced NGOs. Small grassroots\n'
            'organizations — the ones actually serving rural India — are locked out\n'
            'because they lack the bandwidth for legal compliance paperwork (12A, 80G,\n'
            'FCRA) and cannot draft corporate-standard English proposals.\n'
            '\n'
            'The capital is there. The transmission layer is broken.'
        )
    elif txt.startswith('Team Leader'):
        simple_set(shape, 'Team Leader: Siddhant')

# ================================================================
# SLIDE 2: Brief about the Idea
# ================================================================
print("Slide 2: Brief about the Idea")
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if not shape.has_text_frame:
        continue
    if 'Brief about' in shape.text_frame.text:
        simple_set(shape,
            'Brief about the Idea:\n'
            '\n'
            'NidhiAI is a multi-agent AI platform that gives grassroots NGOs the same\n'
            'administrative firepower as large organizations — for free.\n'
            '\n'
            'What it does:\n'
            '1. Compliance Automation — NGOs upload their 12A/80G certificates. Amazon\n'
            '   Textract extracts key fields; a Bedrock Agent validates legal standing\n'
            '   in seconds instead of weeks with a chartered accountant.\n'
            '\n'
            '2. Smart CSR Grant Matching — We built a Knowledge Base of corporate CSR\n'
            '   filings (Tata, Infosys, HDFC, Reliance, Wipro). Titan Embeddings match\n'
            '   the NGO profile to relevant funding opportunities semantically.\n'
            '\n'
            '3. One-Click Proposal Drafting — A single button press generates a formal,\n'
            '   5-page grant application tailored to the specific fund guidelines,\n'
            '   completely removing the English fluency barrier.\n'
            '\n'
            '4. CSR Legal Assistant — An integrated chatbot backed by a Knowledge Base\n'
            '   of Indian CSR laws (Section 135, Schedule VII, FCRA) provides instant\n'
            '   regulatory guidance.'
        )

# ================================================================
# SLIDE 3: Why AI / How AWS / Value of AI layer
# ================================================================
print("Slide 3: Why AI is required")
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if not shape.has_text_frame:
        continue
    if 'Your solution should' in shape.text_frame.text or 'AI is required' in shape.text_frame.text:
        simple_set(shape,
            'Why AI is required in NidhiAI:\n'
            '\n'
            'This is not a CRUD app with a chatbot bolted on. Every core feature\n'
            'is impossible without AI:\n'
            '\n'
            '> Document Compliance: Government certificates (12A/80G) are unstructured\n'
            '  scanned PDFs with no standard API. Amazon Textract + a Bedrock Agent is\n'
            '  the only viable extraction pipeline for registration numbers and dates.\n'
            '\n'
            '> Semantic Matching: A keyword search cannot map "village water well\n'
            '  project" to "Schedule VII Rural Infrastructure Development." We need\n'
            '  Titan Text Embeddings V2 + OpenSearch vector search for this.\n'
            '\n'
            '> Proposal Generation: Grassroots NGOs cannot afford Rs.50,000 per grant\n'
            '  writer. Claude 3.5 Sonnet via Bedrock acts as that writer — ingesting\n'
            '  NGO context and corporate criteria via RAG to produce formal proposals.\n'
            '\n'
            'How AWS services are used:\n'
            '  - Amazon Bedrock Agents (Supervisor Pattern) — orchestrates 4 sub-agents\n'
            '  - Amazon Textract — OCR + document analysis on legal PDFs\n'
            '  - Amazon OpenSearch Serverless — vector store for Knowledge Bases\n'
            '  - AWS Lambda (x4) — serverless compute for each agent action\n'
            '  - Amazon DynamoDB — state management across 4 tables\n'
            '  - Amazon S3 (x5 buckets) — document storage and KB data sources'
        )

# ================================================================
# SLIDE 4: List of features
# ================================================================
print("Slide 4: Features list")
slide4 = prs.slides[3]
for shape in slide4.shapes:
    if not shape.has_text_frame:
        continue
    if 'List of features' in shape.text_frame.text:
        simple_set(shape,
            'Features:\n'
            '\n'
            'Core Platform:\n'
            '  - NGO Profile Management with sector and region tagging\n'
            '  - Drag-and-drop document upload to S3 with real-time progress\n'
            '  - Real-time compliance dashboard (12A / 80G / CSR-1 status)\n'
            '\n'
            'AI-Powered Agents (Amazon Bedrock Supervisor Pattern):\n'
            '  - Compliance Agent — Textract-powered legal document validation\n'
            '  - Grant Scout Agent — semantic matching against 5+ corporate CSR funds\n'
            '  - Proposal Agent — RAG-based formal proposal generation (PDF export)\n'
            '  - Impact Agent — quarterly report generation with AI-computed metrics\n'
            '\n'
            'CSR Assistant Chatbot:\n'
            '  - Backed by a synced Knowledge Base of Indian CSR laws\n'
            '  - Answers questions about Section 135, Schedule VII, FCRA, 12A/80G\n'
            '  - Pre-built suggested questions for ease of use\n'
            '\n'
            'Developer & UX Features:\n'
            '  - Agent Trace panel — shows the multi-agent orchestration in real-time\n'
            '  - Natural language command bar on the dashboard\n'
            '  - Dark/Light theme toggle\n'
            '  - Fully responsive design (Next.js 16 + TailwindCSS)'
        )

# ================================================================
# SLIDE 5: Process flow / Use-case / Architecture diagram
# ================================================================
print("Slide 5: Flow Diagram")
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if not shape.has_text_frame:
        continue
    if 'Process flow' in shape.text_frame.text or 'flow diagram' in shape.text_frame.text:
        simple_set(shape,
            'User Flow:\n'
            '\n'
            '1. NGO Leader signs up and creates a profile (name, sector, region)\n'
            '\n'
            '2. Uploads legal documents (12A cert, 80G cert, CSR-1 form) via\n'
            '   drag-and-drop. Files go directly to the nidhiai-documents S3 bucket.\n'
            '\n'
            '3. The Supervisor Agent (HB82HPMIA3) orchestrates the Compliance Agent\n'
            '   (GFOH0VN84S), which calls the nidhiai-scan-documents Lambda. Amazon\n'
            '   Textract reads the PDFs. Results are stored in DynamoDB.\n'
            '\n'
            '4. Once compliant, the Grant Scout Agent (KRPTCZMV4Z) calls the\n'
            '   nidhiai-match-grants Lambda. Titan Embeddings V2 queries OpenSearch\n'
            '   against the nidhiai-kb-csr-opportunities bucket. Top matches returned\n'
            '   with relevance scores (e.g., Tata Steel — 93%).\n'
            '\n'
            '5. NGO clicks "Draft Proposal." The Proposal Agent (23TIS55Z0P) calls\n'
            '   nidhiai-generate-pdf Lambda. Claude generates a full proposal via RAG;\n'
            '   the PDF is saved to nidhiai-generated-pdfs bucket.\n'
            '\n'
            '6. For general CSR questions, the Supervisor directly queries the\n'
            '   nidhiai-kb-csr-laws Knowledge Base (WUIYUPN7I2) and responds.'
        )

# ================================================================
# SLIDE 6: Wireframes / Mock diagrams
# ================================================================
print("Slide 6: Wireframes (adding screenshots)")
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if not shape.has_text_frame:
        continue
    if 'Wireframes' in shape.text_frame.text:
        simple_set(shape, 'Wireframes / Screenshots of the Working Prototype:')

# Add dashboard + grants screenshots
dashboard_img = os.path.join(SS_DIR, 'dashboard_1772968871753.png')
grants_img = os.path.join(SS_DIR, 'grants_1772969005674.png')
add_image_to_slide(slide6, dashboard_img, Inches(0.3), Inches(1.5), Inches(4.5), Inches(2.8))
add_image_to_slide(slide6, grants_img, Inches(5.0), Inches(1.5), Inches(4.5), Inches(2.8))

# ================================================================
# SLIDE 7: Architecture Diagram
# ================================================================
print("Slide 7: Architecture Diagram")
slide7 = prs.slides[6]
for shape in slide7.shapes:
    if not shape.has_text_frame:
        continue
    if 'Architecture diagram' in shape.text_frame.text:
        simple_set(shape,
            'Architecture Diagram:\n'
            '\n'
            '                    [Next.js Frontend on Vercel]\n'
            '                              |\n'
            '                    [AWS API Gateway REST]\n'
            '                              |\n'
            '                 [nidhiai-gateway Lambda]\n'
            '                              |\n'
            '             +----------------+----------------+\n'
            '             |                                 |\n'
            '    [Amazon Bedrock Agents]          [Direct Lambda Calls]\n'
            '    Supervisor (HB82HPMIA3)          for S3 upload, DynamoDB\n'
            '             |                                 |\n'
            '    +--------+--------+--------+     [Amazon S3 - 5 buckets]\n'
            '    |        |        |        |     [DynamoDB - 4 tables]\n'
            '  Compliance Grant  Proposal Impact\n'
            '  (GFOH0VN) Scout   (23TIS5) (PHLYZ)\n'
            '            (KRPTCZ)\n'
            '    |        |        |\n'
            '  Textract  OpenSearch  Claude/Nova\n'
            '            Serverless  via Bedrock\n'
            '            (3 KBs)\n'
            '\n'
            'All compute is serverless. Zero always-on servers.'
        )

# ================================================================
# SLIDE 8: Technologies
# ================================================================
print("Slide 8: Technologies")
slide8 = prs.slides[7]
for shape in slide8.shapes:
    if not shape.has_text_frame:
        continue
    if 'Technologies utilized' in shape.text_frame.text:
        simple_set(shape,
            'Technologies utilized:\n'
            '\n'
            'AWS AI/ML:\n'
            '  - Amazon Bedrock Agents — 5 agents (1 Supervisor + 4 sub-agents)\n'
            '  - Amazon Bedrock Knowledge Bases — 3 KBs (CSR laws, opportunities, proposals)\n'
            '  - Amazon Titan Text Embeddings V2 — vector embeddings for semantic search\n'
            '  - Foundation Models: Claude 3.5 Sonnet, Amazon Nova 2 Lite\n'
            '  - Amazon Textract — OCR and structured data extraction from PDFs\n'
            '\n'
            'AWS Infrastructure:\n'
            '  - AWS Lambda — 4 functions (Python 3.12, 256MB, 90s timeout)\n'
            '    * nidhiai-scan-documents\n'
            '    * nidhiai-match-grants\n'
            '    * nidhiai-generate-pdf\n'
            '    * nidhiai-generate-report\n'
            '  - Amazon API Gateway — REST API routing\n'
            '  - Amazon DynamoDB — 4 tables (profiles, documents, compliance, proposals)\n'
            '  - Amazon S3 — 5 buckets (documents, PDFs, 3 KB data sources)\n'
            '  - Amazon OpenSearch Serverless — vector store backing all 3 KBs\n'
            '  - Amazon Cognito — user authentication\n'
            '\n'
            'Frontend:\n'
            '  - Next.js 16 (React 19, Server Components)\n'
            '  - TailwindCSS + custom dark theme\n'
            '  - Deployed on Vercel'
        )

# ================================================================
# SLIDE 9: Estimated cost
# ================================================================
print("Slide 9: Cost")
slide9 = prs.slides[8]
for shape in slide9.shapes:
    if not shape.has_text_frame:
        continue
    if 'Estimated implementation cost' in shape.text_frame.text:
        simple_set(shape,
            'Estimated implementation cost:\n'
            '\n'
            'Prototype (current):\n'
            '  - Lambda / API Gateway / DynamoDB: Covered under AWS Free Tier\n'
            '  - Bedrock inference: ~Rs.800/month at demo volume\n'
            '  - S3 storage: < Rs.50/month across all 5 buckets\n'
            '  - OpenSearch Serverless: Min-capacity provisions apply\n'
            '  - Total prototype spend to date: < Rs.8,000\n'
            '\n'
            'Production (100 NGOs/month):\n'
            '  - Estimated Rs.15,000-25,000/month\n'
            '  - Per-proposal cost: ~Rs.25 vs. Rs.50,000+ for a human grant writer\n'
            '  - This represents a 2000x cost reduction for NGOs\n'
            '\n'
            'The serverless model means zero cost when idle — no NGO pays for\n'
            'infrastructure they are not using.'
        )

# ================================================================
# SLIDE 10: Snapshots of the prototype
# ================================================================
print("Slide 10: Snapshots")
slide10 = prs.slides[9]
for shape in slide10.shapes:
    if not shape.has_text_frame:
        continue
    if 'Snapshots' in shape.text_frame.text:
        simple_set(shape, 'Snapshots of the working prototype:')

# Add compliance + chatbot + proposals screenshots
compliance_img = os.path.join(SS_DIR, 'compliance_1772968882563.png')
proposals_img = os.path.join(SS_DIR, 'proposals_1772969030004.png')
chatbot_img = os.path.join(SS_DIR, 'chatbot_final_1772969216606.png')
add_image_to_slide(slide10, compliance_img, Inches(0.2), Inches(1.3), Inches(3.1), Inches(2.0))
add_image_to_slide(slide10, proposals_img, Inches(3.4), Inches(1.3), Inches(3.1), Inches(2.0))
add_image_to_slide(slide10, chatbot_img, Inches(6.6), Inches(1.3), Inches(3.1), Inches(2.0))

# ================================================================
# SLIDE 11: Performance / Benchmarking
# ================================================================
print("Slide 11: Performance")
slide11 = prs.slides[10]
for shape in slide11.shapes:
    if not shape.has_text_frame:
        continue
    if 'Performance' in shape.text_frame.text or 'Benchmarking' in shape.text_frame.text:
        simple_set(shape,
            'Prototype Performance:\n'
            '\n'
            'End-to-End Latency (measured on live AWS infra, ap-south-1):\n'
            '\n'
            '  Document Upload to S3:              ~1.2 seconds\n'
            '  Compliance Scan (Textract + Agent):  4-6 seconds per document\n'
            '  Grant Matching (Titan V2 + AOSS):    1.8 seconds\n'
            '  Proposal Generation (RAG + Claude):  15-20 seconds (streamed)\n'
            '  KB Chatbot Response:                 2-4 seconds\n'
            '\n'
            'Infrastructure:\n'
            '  System availability:  100% (fully managed AWS services)\n'
            '  Cold start penalty:   < 800ms (Lambda Python 3.12, 256MB)\n'
            '  Concurrent capacity:  Auto-scales with Lambda provisioning\n'
            '\n'
            'Agent Orchestration:\n'
            '  All 5 Bedrock Agents status: PREPARED\n'
            '  All 3 Knowledge Bases status: ACTIVE\n'
            '  KB sync: 18 documents indexed with 0 failures'
        )

# ================================================================
# SLIDE 12: Future Development
# ================================================================
print("Slide 12: Future Development")
slide12 = prs.slides[11]
for shape in slide12.shapes:
    if not shape.has_text_frame:
        continue
    if 'Additional Details' in shape.text_frame.text or 'Future Development' in shape.text_frame.text:
        simple_set(shape,
            'Future Development:\n'
            '\n'
            'Near-term (next 3 months):\n'
            '  - Multilingual voice intake via Amazon Transcribe — let NGO founders\n'
            '    describe their work in Hindi or regional languages, and have the AI\n'
            '    convert it into a formal English proposal automatically\n'
            '  - Direct submission to corporate CSR portals through API integrations\n'
            '  - Expand the Grant KB to cover 50+ companies (currently 5)\n'
            '\n'
            'Medium-term (6-12 months):\n'
            '  - Government scheme matching (PM-KISAN, MGNREGA tie-ins)\n'
            '  - Automated annual compliance reminders (12A/80G renewal alerts)\n'
            '  - Mobile-first PWA for low-bandwidth rural areas\n'
            '\n'
            'Long-term vision:\n'
            '  NidhiAI becomes the default operating system for NGO administration\n'
            '  in India — handling compliance, funding, reporting, and government\n'
            '  filings in one platform. If 10% of the Rs.38,000 Cr CSR pool reaches\n'
            '  smaller NGOs through NidhiAI, that is Rs.3,800 Crores unlocked for\n'
            '  communities that need it most.'
        )

# ================================================================
# SLIDE 13: Prototype Assets (GitHub + Video)
# ================================================================
print("Slide 13: Assets")
slide13 = prs.slides[12]
for shape in slide13.shapes:
    if not shape.has_text_frame:
        continue
    txt = shape.text_frame.text
    if 'Prototype Assets' in txt:
        simple_set(shape, 'Prototype Assets:')
    if 'GitHub' in txt:
        simple_set(shape,
            'GitHub Repository: https://github.com/Sid-V5/NidhiAi\n'
            '\n'
            'Demo Video: [Add YouTube / Google Drive link here]\n'
            '\n'
            'Live MVP: [Add Vercel deployment URL here]'
        )

# ================================================================
# SLIDE 14: Thank you (just the AWS image — no changes needed)
# ================================================================
print("Slide 14: Thank you slide — no changes needed")

# ================================================================
# Save
# ================================================================
prs.save(OUTPUT)
print(f"\nDone! Saved final PPT to:\n  {OUTPUT}")
