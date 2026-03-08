"""
NidhiAI - AI for Bharat Hackathon Presentation Generator
Generates a professional PPTX following the official template structure.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── COLORS ───
DARK_BG = RGBColor(0x0D, 0x1B, 0x2A)       # Deep navy
ACCENT_ORANGE = RGBColor(0xFF, 0x99, 0x33)  # AWS Orange
ACCENT_TEAL = RGBColor(0x00, 0xD4, 0xAA)    # Bright teal
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
RED = RGBColor(0xFF, 0x44, 0x44)
GREEN = RGBColor(0x00, 0xCC, 0x66)
LIGHT_BG = RGBColor(0x14, 0x24, 0x3B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_box(slide, left, top, width, height, items, font_size=13, color=WHITE, bullet_color=ACCENT_TEAL):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.alignment = PP_ALIGN.JUSTIFY
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color, text="", font_size=11, font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(2)
    return shape

def add_line_connector(slide, x1, y1, x2, y2, color=ACCENT_TEAL, width=2):
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector

def add_heading_bar(slide, text):
    """Adds a consistent heading bar at the top of each slide."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = LIGHT_BG
    bar.line.fill.background()
    add_text_box(slide, 0.5, 0.2, 12, 0.6, text, font_size=26, color=ACCENT_ORANGE, bold=True)


# ═══════════════════════════════════════════════════
# SLIDE 1: TITLE SLIDE
# ═══════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1)

# Title
add_text_box(slide1, 1, 1.0, 11, 1.2, "NidhiAI (निधिAI)", font_size=48, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide1, 1.5, 2.3, 10, 0.8, "Unlocking 38,000 Crore of CSR Capital for Bharat's Changemakers", font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)

# Divider line
div = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.3), Inches(5.333), Inches(0.03))
div.fill.solid()
div.fill.fore_color.rgb = ACCENT_TEAL
div.line.fill.background()

# Team info
add_text_box(slide1, 1, 3.8, 5.5, 0.4, "Team Name:  [YOUR TEAM NAME]", font_size=16, color=LIGHT_GRAY)
add_text_box(slide1, 1, 4.3, 5.5, 0.4, "Team Leader:  [YOUR NAME]", font_size=16, color=LIGHT_GRAY)

# Problem Statement Field (Combined Track + Summary)
problem_text = "Problem Statement: Track 03 - AI for Communities\nSpecific Problem: The 'Last Mile' CSR Funding Gap (90% Misallocation)"
add_text_box(slide1, 7, 3.8, 5.8, 1.2, problem_text, font_size=16, color=ACCENT_TEAL, alignment=PP_ALIGN.RIGHT)

# Bottom bar
bottom = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7))
bottom.fill.solid()
bottom.fill.fore_color.rgb = LIGHT_BG
bottom.line.fill.background()
add_text_box(slide1, 0.5, 6.85, 12, 0.5, "AI for Bharat Hackathon 2026  |  Powered by AWS  |  Amazon Bedrock • Amazon Q • Amazon Textract", font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 2: BRIEF ABOUT THE IDEA
# ═══════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)
add_heading_bar(slide2, "Brief About the Idea")

# Left column - Problem
add_text_box(slide2, 0.5, 1.2, 6, 0.5, "THE PROBLEM", font_size=18, color=RED, bold=True)
add_bullet_box(slide2, 0.5, 1.8, 6, 3.5, [
    "India's companies spend 18,000+ Cr/year on CSR (Section 135), projected to reach 38,000 Cr (CSRBOX ICOR 2025)",
    "But 70% of funds go to the same top 500 companies' preferred NGOs (give.do)",
    "3 Million+ civil society organizations are underserved (CBI/ICNL)",
    "Triple Barrier: Information Gap + Compliance Complexity + English Proposal Barrier",
    "A village schoolteacher feeding 200 children cannot write a corporate grant proposal",
], font_size=14)

# Right column - Solution
add_text_box(slide2, 7, 1.2, 6, 0.5, "THE SOLUTION", font_size=18, color=GREEN, bold=True)
add_bullet_box(slide2, 7, 1.8, 5.8, 3.5, [
    "NidhiAI = AI-powered Fundraising Officer for small NGOs",
    "Upload documents > Check compliance > Find grants > Write proposals",
    "Uses AWS Bedrock Multi-Agent Supervisor Pattern",
    "4 specialized AI agents working as a coordinated team",
    "Turns a 30-day manual process into a 30-minute automated one",
], font_size=14)

# Bottom stat bar
stat_bar = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5), Inches(12.333), Inches(1.2))
stat_bar.fill.solid()
stat_bar.fill.fore_color.rgb = LIGHT_BG
stat_bar.line.fill.background()
add_text_box(slide2, 1, 5.6, 3.5, 0.9, "38,000 Cr\nCSR Market (Projected)", font_size=20, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide2, 4.5, 5.6, 3.5, 0.9, "3 Million+\nCivil Society Orgs", font_size=20, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide2, 8.5, 5.6, 4, 0.9, "70% Funds\nConcentrated in Top NGOs", font_size=20, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide2, 0.5, 6.75, 12.333, 0.3, "Sources: CSRBOX ICOR 2025 Report | give.do CSR Analysis | CBI/ICNL India Registry", font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 3: HOW DIFFERENT / HOW IT SOLVES / USP
# ═══════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)
add_heading_bar(slide3, "Solution Differentiators and USP")

# 3 boxes answering the 3 template questions
qa_data = [
    ("How is it different\nfrom existing ideas?", "No tool automates the FULL\nNGO fundraising lifecycle:\ncompliance + discovery +\nproposal + impact tracking.\nExisting platforms are siloed.\nNidhiAI is an autonomous\nmulti-agent system that\nDOES the work end-to-end."),
    ("How will it\nsolve the problem?", "4 AI Agents work together:\n1. Compliance Agent scans docs\n2. Scout Agent finds grants\n3. Proposal Agent writes PDFs\n4. Impact Agent tracks results\nAll orchestrated by a Bedrock\nSupervisor Agent automatically."),
    ("USP of the\nProposed Solution", "'Sales Intelligence for Charity'\n- B2B prospecting for NGOs\n- Deep Legal RAG on MCA data\n  (zero hallucination)\n- Bedrock Supervisor Pattern\n  (not a chatbot, a team of\n  specialized AI workers)"),
]

for i, (title, desc) in enumerate(qa_data):
    x = 0.5 + i * 4.2
    box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(3.8), Inches(5.0))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = ACCENT_TEAL
    box.line.width = Pt(1.5)
    add_text_box(slide3, x + 0.2, 1.5, 3.4, 0.9, title, font_size=16, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide3, x + 0.3, 2.6, 3.2, 3.2, desc, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide3, 0.5, 6.5, 12, 0.6, "We don't just give advice. We do the work. NidhiAI automates the ENTIRE fundraising lifecycle.", font_size=15, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 4: LIST OF FEATURES
# ═══════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)
add_heading_bar(slide4, "Key Features of NidhiAI")

features = [
    ("Smart Compliance Check", "Upload 12A, 80G, CSR-1 certs. Textract scans and validates. Instant red/green status indicators."),
    ("AI Grant Discovery", "Semantic search across 100+ company CSR reports. Ranked matches based on NGO cause and location."),
    ("Auto Proposal Writer", "Generates professional 5-page PDF proposals tailored to each donor's priorities and language."),
    ("Impact Report Generator", "Auto-generates quarterly impact reports from field data to maintain donor relationships."),
    ("CSR Helpdesk (Amazon Q)", "Embedded chatbot answering questions about CSR laws, FCRA, Schedule VII in plain language."),
    ("Enterprise Security", "S3 encryption at rest, TLS in transit, Cognito auth, data isolation per NGO."),
]

for i, (title, desc) in enumerate(features):
    row = i // 2
    col = i % 2
    x = 0.5 + col * 6.4
    y = 1.3 + row * 1.8

    box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(6.0), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.fill.background()
    add_text_box(slide4, x + 0.2, y + 0.1, 5.6, 0.5, title, font_size=15, color=ACCENT_ORANGE, bold=True)
    add_text_box(slide4, x + 0.2, y + 0.6, 5.6, 0.7, desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.JUSTIFY)


# ═══════════════════════════════════════════════════
# SLIDE 5: PROCESS FLOW DIAGRAM
# ═══════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)
add_heading_bar(slide5, "Process Flow: How NidhiAI Works")

# Flow: 5 steps left to right
steps = [
    ("👤 NGO Leader\nUploads Docs", ACCENT_TEAL),
    ("⚖️ Compliance\nAgent Scans", RGBColor(0x44, 0x88, 0xFF)),
    ("🕵️ Scout Agent\nFinds Grants", RGBColor(0xFF, 0x88, 0x00)),
    ("📝 Proposal Agent\nWrites PDF", RGBColor(0xAA, 0x44, 0xFF)),
    ("📄 NGO Gets\nReady Proposal", GREEN),
]

step_width = 2.0
gap = 0.6
total_width = len(steps) * step_width + (len(steps) - 1) * gap
start_x = (13.333 - total_width) / 2

for i, (label, color) in enumerate(steps):
    x = start_x + i * (step_width + gap)
    y = 2.2

    box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(step_width), Inches(1.8))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    if i < len(steps) - 1:
        arrow_x = x + step_width
        arrow = slide5.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(arrow_x + 0.05), Inches(y + 0.65), Inches(gap - 0.1), Inches(0.5))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_TEAL
        arrow.line.fill.background()

# Supervisor label
add_rounded_rect(slide5, 3, 4.5, 7.333, 0.7, LIGHT_BG, "🤖 All Orchestrated by the SUPERVISOR AGENT (Bedrock Supervisor Pattern)", font_size=14, font_color=ACCENT_ORANGE)

# Detail boxes below
detail_data = [
    ("Step 1", "Textract extracts\ndate fields from\nscanned certificates"),
    ("Step 2", "Bedrock KB semantic\nsearch matches NGO\nto CSR opportunities"),
    ("Step 3", "Claude 3.5 Sonnet\ngenerates tailored\nproposal as PDF"),
]
for i, (step, detail) in enumerate(detail_data):
    x = 1.5 + i * 4.0
    add_text_box(slide5, x, 5.5, 3.5, 0.3, step, font_size=12, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide5, x, 5.8, 3.5, 0.9, detail, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 6: WIREFRAMES / MOCKUP
# ═══════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6)
add_heading_bar(slide6, "Wireframes / Mock Diagrams (UI Preview)")

# Dashboard mockup - left panel
dash_bg = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(6), Inches(5.5))
dash_bg.fill.solid()
dash_bg.fill.fore_color.rgb = LIGHT_BG
dash_bg.line.color.rgb = ACCENT_TEAL
dash_bg.line.width = Pt(1)

add_text_box(slide6, 0.7, 1.5, 5.5, 0.4, "NidhiAI Dashboard", font_size=16, color=ACCENT_ORANGE, bold=True)
# Nav items
nav_items = ["📋 Dashboard", "📤 Upload Documents", "🔍 Find Grants", "📝 Generate Proposal", "📊 Impact Reports", "💬 CSR Helpdesk"]
for i, item in enumerate(nav_items):
    btn = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.1 + i * 0.65), Inches(5.2), Inches(0.5))
    btn.fill.solid()
    btn.fill.fore_color.rgb = RGBColor(0x1A, 0x2E, 0x4A) if i > 0 else ACCENT_TEAL
    btn.line.fill.background()
    tf = btn.text_frame
    p = tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"

# Right panel - compliance check result
right_bg = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.3), Inches(5.8), Inches(5.5))
right_bg.fill.solid()
right_bg.fill.fore_color.rgb = LIGHT_BG
right_bg.line.color.rgb = ACCENT_TEAL
right_bg.line.width = Pt(1)

add_text_box(slide6, 7.2, 1.5, 5.3, 0.4, "Compliance Check Results", font_size=16, color=ACCENT_ORANGE, bold=True)

checks = [
    ("✅  12A Registration", "Valid till: Dec 2028", GREEN),
    ("✅  80G Certificate", "Valid till: Nov 2027", GREEN),
    ("CSR-1 Filing", "EXPIRED: Jan 2025 - Please Renew", RED),
    ("✅  PAN Verification", "Verified: AAATX1234F", GREEN),
]
for i, (label, detail, color) in enumerate(checks):
    y = 2.2 + i * 0.95
    chk = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(y), Inches(5.2), Inches(0.8))
    chk.fill.solid()
    chk.fill.fore_color.rgb = RGBColor(0x1A, 0x2E, 0x4A)
    chk.line.fill.background()
    add_text_box(slide6, 7.5, y + 0.05, 4.8, 0.35, label, font_size=13, color=color, bold=True)
    add_text_box(slide6, 7.5, y + 0.38, 4.8, 0.35, detail, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════
# SLIDE 7: ARCHITECTURE DIAGRAM
# ═══════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7)
add_heading_bar(slide7, "Architecture Diagram: AWS Services")

# --- TOP ROW: User Interface Layer ---
add_rounded_rect(slide7, 0.3, 1.3, 2.2, 0.9, ACCENT_TEAL, "NGO User\n(Web Dashboard)", font_size=11)
add_rounded_rect(slide7, 2.8, 1.3, 2.2, 0.9, RGBColor(0x88, 0x44, 0xCC), "Amazon Q Business\n(CSR FAQ Chat)", font_size=10)
add_rounded_rect(slide7, 5.3, 1.3, 2.4, 0.9, RGBColor(0x44, 0x66, 0xBB), "API Gateway\n+ Amazon Cognito", font_size=10)

# connectors: user -> Q, user -> API GW
add_line_connector(slide7, 2.5, 1.75, 2.8, 1.75)
add_line_connector(slide7, 5.0, 1.75, 5.3, 1.75)

# --- MIDDLE ROW: Supervisor Agent ---
add_rounded_rect(slide7, 3.5, 2.6, 5.5, 1.0, ACCENT_ORANGE, "BEDROCK SUPERVISOR AGENT (Claude 3.5 Sonnet)", font_size=14)
# Connector: API GW -> Supervisor
add_line_connector(slide7, 6.5, 2.2, 6.5, 2.6)

# Episodic Memory (right of supervisor)
add_rounded_rect(slide7, 9.5, 2.6, 2.2, 1.0, RGBColor(0x55, 0x55, 0x88), "Episodic Memory\n(Session Context)", font_size=10)
add_line_connector(slide7, 9.0, 3.1, 9.5, 3.1)

# --- BOTTOM ROW: 4 Specialist Agents ---
agents = [
    ("Compliance\nAgent", RGBColor(0x44, 0x88, 0xFF)),
    ("Grant Scout\nAgent", RGBColor(0xFF, 0x88, 0x00)),
    ("Proposal\nAgent", RGBColor(0xAA, 0x44, 0xFF)),
    ("Impact\nAgent", RGBColor(0x00, 0xCC, 0x66)),
]
for i, (name, color) in enumerate(agents):
    x = 0.3 + i * 3.2
    add_rounded_rect(slide7, x, 4.1, 2.7, 0.9, color, name, font_size=11)
    # Connector: Supervisor -> Agent
    cx = x + 1.35
    add_line_connector(slide7, cx, 3.6, cx, 4.1)

# --- RIGHT COLUMN: AWS Services (stacked) ---
aws_svc = [
    ("Amazon Textract\n(Doc Scanning)", RGBColor(0x22, 0x55, 0x99)),
    ("KB: CSR Laws\n(MCA/Schedule VII)", RGBColor(0x22, 0x55, 0x99)),
    ("KB: CSR Opportunities\n(Company CSR Reports)", RGBColor(0x22, 0x55, 0x99)),
    ("KB: Winning Proposals\n(Template Library)", RGBColor(0x22, 0x55, 0x99)),
]
for i, (svc, color) in enumerate(aws_svc):
    y = 4.1 + i * 0.7
    # Stagger to show which agent connects
    if i < 4:
        add_rounded_rect(slide7, 10.5, y - 0.6 + i * 0.55, 2.5, 0.5, color, svc, font_size=8)

# Actually let's place them cleaner - under the agents
# Re-do the right column as a proper service stack
# Clear and re-draw the right column
add_rounded_rect(slide7, 10.2, 1.3, 2.8, 0.7, RGBColor(0x22, 0x55, 0x99), "Amazon Textract", font_size=10)
add_rounded_rect(slide7, 10.2, 2.1, 2.8, 0.7, RGBColor(0x22, 0x55, 0x99), "KB: CSR Laws", font_size=10)
add_rounded_rect(slide7, 10.2, 2.9, 2.8, 0.7, RGBColor(0x22, 0x55, 0x99), "KB: Opportunities", font_size=10)
add_rounded_rect(slide7, 10.2, 3.7, 2.8, 0.7, RGBColor(0x22, 0x55, 0x99), "KB: Proposals", font_size=10)

# Bottom services row
bottom_svc = [
    ("Lambda: Doc Scanner", RGBColor(0x22, 0x55, 0x99)),
    ("Lambda: Grant Matcher", RGBColor(0x22, 0x55, 0x99)),
    ("Lambda: PDF Generator", RGBColor(0x22, 0x55, 0x99)),
    ("Amazon S3 (Storage)", RGBColor(0x22, 0x55, 0x99)),
    ("OpenSearch Serverless", RGBColor(0x22, 0x55, 0x99)),
]
for i, (svc, color) in enumerate(bottom_svc):
    x = 0.3 + i * 2.6
    add_rounded_rect(slide7, x, 5.5, 2.3, 0.6, color, svc, font_size=9)

# Bottom label
add_text_box(slide7, 0.5, 6.4, 12, 0.5, "All services in AWS ap-south-1 (Mumbai) | Serverless Architecture | 10 AWS Services Total", font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 8: TECHNOLOGIES
# ═══════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8)
add_heading_bar(slide8, "Technologies Used in the Solution")

tech_categories = [
    ("AI & ML Layer", [
        "Amazon Bedrock Agents (Supervisor Mode)",
        "Claude 3.5 Sonnet (Foundation Model)",
        "Amazon Bedrock Knowledge Bases (x3)",
        "Amazon Titan Text Embeddings v2",
    ]),
    ("Document Intelligence", [
        "Amazon Textract (OCR + Form Extraction)",
        "Scans 12A, 80G, CSR-1 Certificates",
        "94% accuracy on structured forms",
        "Available in Mumbai (ap-south-1)",
    ]),
    ("Compute & Storage", [
        "AWS Lambda (Python 3.12) - 3 Functions",
        "Amazon S3 (Document + PDF Storage)",
        "Amazon OpenSearch Serverless (Vector DB)",
        "Amazon API Gateway (REST APIs)",
    ]),
    ("Frontend & Auth", [
        "Next.js (React) - Web Dashboard",
        "Amazon Cognito (User Auth)",
        "Amazon Q Business (Embedded CSR FAQ)",
        "Amazon Q Developer (Dev Workflow)",
    ]),
]

for i, (category, items) in enumerate(tech_categories):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.5
    y = 1.3 + row * 3.0

    box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(6.0), Inches(2.7))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = ACCENT_TEAL
    box.line.width = Pt(1)

    add_text_box(slide8, x + 0.2, y + 0.1, 5.5, 0.4, category, font_size=16, color=ACCENT_ORANGE, bold=True)
    add_bullet_box(slide8, x + 0.2, y + 0.6, 5.5, 1.8, items, font_size=12)


# ═══════════════════════════════════════════════════
# SLIDE 9: BUSINESS MODEL & UNIT ECONOMICS
# ═══════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9)
add_heading_bar(slide9, "Business Model and Unit Economics")

# LEFT COLUMN: Unit Economics comparison
left_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5))
left_box.fill.solid()
left_box.fill.fore_color.rgb = LIGHT_BG
left_box.line.color.rgb = ACCENT_TEAL
left_box.line.width = Pt(1)

add_text_box(slide9, 0.7, 1.4, 5.5, 0.4, "Unit Economics: NidhiAI vs Traditional", font_size=16, color=ACCENT_ORANGE, bold=True)

# Comparison rows
comparisons = [
    ("Metric", "Traditional", "NidhiAI", ACCENT_ORANGE),
    ("Cost per proposal", "Rs.50,000+", "< Rs.10", WHITE),
    ("Time per proposal", "2-4 weeks", "< 10 minutes", WHITE),
    ("Grant discovery", "Manual search", "AI-matched (100+ cos)", WHITE),
    ("Compliance check", "CA fees Rs.5,000+", "Automated via Textract", WHITE),
    ("Scaling to 100 NGOs", "50 grant writers", "Same AI system", WHITE),
]

for i, (metric, trad, nidhi, color) in enumerate(comparisons):
    y = 2.0 + i * 0.55
    bg = RGBColor(0x1A, 0x3A, 0x5A) if i == 0 else (DARK_BG if i % 2 == 0 else LIGHT_BG)
    row = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y), Inches(5.5), Inches(0.5))
    row.fill.solid()
    row.fill.fore_color.rgb = bg
    row.line.fill.background()
    bold = i == 0
    add_text_box(slide9, 0.8, y + 0.05, 2.0, 0.4, metric, font_size=12, color=color, bold=bold)
    add_text_box(slide9, 2.9, y + 0.05, 1.6, 0.4, trad, font_size=12, color=RGBColor(0xFF, 0x66, 0x66) if i > 0 else color, bold=bold, alignment=PP_ALIGN.CENTER)
    add_text_box(slide9, 4.5, y + 0.05, 1.6, 0.4, nidhi, font_size=12, color=GREEN if i > 0 else color, bold=bold, alignment=PP_ALIGN.CENTER)

# Cost savings callout
savings_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.4), Inches(5.5), Inches(1.0))
savings_box.fill.solid()
savings_box.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x1A)
savings_box.line.color.rgb = GREEN
savings_box.line.width = Pt(1)
add_text_box(slide9, 0.9, 5.5, 5.1, 0.7, "5000x cheaper per proposal | 200x faster delivery\nHackathon Phase: Rs.0 (AWS-provisioned credits)", font_size=13, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# RIGHT COLUMN: Revenue Model
right_box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5))
right_box.fill.solid()
right_box.fill.fore_color.rgb = LIGHT_BG
right_box.line.color.rgb = ACCENT_ORANGE
right_box.line.width = Pt(1)

add_text_box(slide9, 7.2, 1.4, 5.3, 0.4, "Revenue Model (100 NGOs)", font_size=16, color=ACCENT_ORANGE, bold=True)

revenue_rows = [
    ("Metric", "Value", ACCENT_ORANGE),
    ("Free Tier", "5 proposals/month (Rs.0)", WHITE),
    ("Premium Plan", "Unlimited (Rs.999/NGO/month)", WHITE),
    ("If 100 NGOs on Premium", "Rs.99,900/month revenue", WHITE),
    ("AWS Infra Cost (100 NGOs)", "~Rs.17,000/month (~$204)", WHITE),
    ("Gross Margin", "~83%", GREEN),
    ("Break-even", "18 paying NGOs", WHITE),
    ("B2B: CSR Consulting Firms", "Rs.9,999/month enterprise", WHITE),
    ("Target Market", "3M+ civil society orgs (CBI/ICNL)", WHITE),
]

for i, (metric, value, color) in enumerate(revenue_rows):
    y = 2.0 + i * 0.5
    bg = RGBColor(0x1A, 0x3A, 0x5A) if i == 0 else (DARK_BG if i % 2 == 0 else LIGHT_BG)
    row = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(y), Inches(5.3), Inches(0.45))
    row.fill.solid()
    row.fill.fore_color.rgb = bg
    row.line.fill.background()
    bold = i == 0 or i == 5
    add_text_box(slide9, 7.3, y + 0.03, 2.2, 0.35, metric, font_size=12, color=color, bold=bold)
    add_text_box(slide9, 9.6, y + 0.03, 2.7, 0.35, value, font_size=12, color=color, bold=bold, alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════
# SLIDE 10: ADDITIONAL REQUIREMENTS + IMPACT
# ═══════════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10)
add_heading_bar(slide10, "Impact for India and Additional Details")

# Left: Impact
impact_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5))
impact_box.fill.solid()
impact_box.fill.fore_color.rgb = LIGHT_BG
impact_box.line.color.rgb = ACCENT_TEAL
impact_box.line.width = Pt(1)

add_text_box(slide10, 0.7, 1.5, 5.5, 0.5, "Impact for India", font_size=18, color=ACCENT_ORANGE, bold=True)
add_bullet_box(slide10, 0.7, 2.1, 5.5, 4.5, [
    "Economic: Unlocks 18,000+ Cr/year CSR market (projected 38,000 Cr) for 3M+ orgs (CSRBOX)",
    "If 1,000 NGOs get 5 Lakh each = 50 Crore moved to grassroots",
    "Social: Directs funds to remote tribal areas and aspirational districts",
    "Digital India: Incentivizes NGOs to digitize records and adopt tech",
    "Efficiency: Reduces fundraising cost from ~20% to less than 1%",
    "Frees NGO leaders to focus on impact, not paperwork",
    "Pilot Plan: 5 NGOs in Jharkhand tribal belt in Q1 2026",
], font_size=12)

# Right: AWS Services Justification + Kiro + Additional
add_box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.3), Inches(5.8), Inches(5.5))
add_box.fill.solid()
add_box.fill.fore_color.rgb = LIGHT_BG
add_box.line.color.rgb = ACCENT_TEAL
add_box.line.width = Pt(1)

add_text_box(slide10, 7.2, 1.5, 5.3, 0.5, "AWS Services and Hackathon Details", font_size=18, color=ACCENT_ORANGE, bold=True)
add_bullet_box(slide10, 7.2, 2.1, 5.3, 4.5, [
    "Amazon Bedrock: Supervisor + 4 Agents + Knowledge Bases",
    "Amazon Q Business: Embedded CSR FAQ chatbot for users",
    "Amazon Q Developer: Used for building the prototype",
    "Amazon Textract: Document scanning for compliance checks",
    "Kiro: Used for spec-driven development (requirements.md, design.md)",
    "Region: ap-south-1 (Mumbai) for low latency in India",
    "Scalability: 100% serverless, pay-per-use pricing",
    "Future: Multilingual support (Hindi, Marathi) + WhatsApp bot",
], font_size=12)


# ═══════════════════════════════════════════════════
# SLIDE 11: THANK YOU / CTA
# ═══════════════════════════════════════════════════
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide11)

add_text_box(slide11, 1, 1.5, 11, 1.0, "NidhiAI (निधिAI)", font_size=48, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide11, 1.5, 2.8, 10, 1.0, "Let's Fund the Changemakers of Bharat.", font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)

# Divider
div2 = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.0), Inches(5.333), Inches(0.03))
div2.fill.solid()
div2.fill.fore_color.rgb = ACCENT_TEAL
div2.line.fill.background()

add_text_box(slide11, 1, 4.5, 11, 0.5, "Team: [YOUR TEAM NAME]", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide11, 1, 5.1, 11, 0.5, "Track 03: AI for Communities, Access & Public Impact", font_size=16, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
add_text_box(slide11, 1, 5.7, 11, 0.5, "Powered by AWS Bedrock  •  Amazon Q  •  Amazon Textract", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom bar
bottom2 = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7))
bottom2.fill.solid()
bottom2.fill.fore_color.rgb = LIGHT_BG
bottom2.line.fill.background()
add_text_box(slide11, 0.5, 6.85, 12, 0.5, "AI for Bharat Hackathon 2026  |  Thank You", font_size=14, color=ACCENT_ORANGE, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════
output_path = r'C:\Users\Siddhant\test2\NidhiAI_Idea_Submission.pptx'
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
