"""
NidhiAI - Prototype Development Submission PPT Generator
Matches the Idea Submission visual style: dark theme, orange/teal accents,
rounded rectangle cards, formatted tables.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── COLORS ───
DARK_BG = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT_ORANGE = RGBColor(0xFF, 0x99, 0x33)
ACCENT_TEAL = RGBColor(0x00, 0xD4, 0xAA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
RED = RGBColor(0xFF, 0x44, 0x44)
GREEN = RGBColor(0x00, 0xCC, 0x66)
LIGHT_BG = RGBColor(0x14, 0x24, 0x3B)
BLUE = RGBColor(0x44, 0x88, 0xFF)
PURPLE = RGBColor(0xAA, 0x44, 0xFF)
YELLOW = RGBColor(0xFF, 0xCC, 0x00)
DIM = RGBColor(0x88, 0x99, 0xAA)

SS_DIR = r'C:\Users\Siddhant\.gemini\antigravity\brain\ffadbb09-a544-4a00-a7d0-c6df051cd0f2'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ─── HELPERS ───
def add_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def tb(slide, left, top, width, height, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box

def bullets(slide, left, top, width, height, items, sz=12, color=WHITE):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(4)
    return box

def card(slide, left, top, width, height, fill=LIGHT_BG, border=None, text="", sz=11, fc=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = fc
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    return shape

def heading(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = LIGHT_BG
    bar.line.fill.background()
    tb(slide, 0.5, 0.2, 12, 0.6, text, sz=26, color=ACCENT_ORANGE, bold=True)

def table_row(slide, x, y, w, metric, value, i, mc=WHITE, vc=WHITE, is_header=False):
    bg = RGBColor(0x1A, 0x3A, 0x5A) if is_header else (DARK_BG if i % 2 == 0 else LIGHT_BG)
    row = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.45))
    row.fill.solid()
    row.fill.fore_color.rgb = bg
    row.line.fill.background()
    tb(slide, x + 0.15, y + 0.03, w * 0.5, 0.35, metric, sz=11, color=mc, bold=is_header)
    tb(slide, x + w * 0.5, y + 0.03, w * 0.45, 0.35, value, sz=11, color=vc, bold=is_header, align=PP_ALIGN.RIGHT)

def footer(slide, text="AI for Bharat Hackathon 2026  |  Powered by AWS  |  Track 03"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.9), Inches(13.333), Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = LIGHT_BG
    bar.line.fill.background()
    tb(slide, 0.5, 6.95, 12, 0.4, text, sz=11, color=DIM, align=PP_ALIGN.CENTER)

def add_img(slide, filename, left, top, width, height=None):
    path = os.path.join(SS_DIR, filename)
    if os.path.exists(path):
        if height:
            slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
        else:
            slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width))
        return True
    print(f"  WARN: {filename} not found")
    return False


# ═══════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tb(s, 1, 1.0, 11, 1.2, "NidhiAI (निधिAI)", sz=48, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
tb(s, 1.5, 2.3, 10, 0.8, "Prototype Development Submission", sz=24, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, 1.5, 3.1, 10, 0.6, "Unlocking Rs.38,000 Crore of CSR Capital for Bharat's Changemakers", sz=16, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

# Divider
div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.0), Inches(5.333), Inches(0.03))
div.fill.solid(); div.fill.fore_color.rgb = ACCENT_TEAL; div.line.fill.background()

tb(s, 1, 4.5, 5.5, 0.4, "Team Name:  NidhiAI", sz=16, color=LIGHT_GRAY)
tb(s, 1, 5.0, 5.5, 0.4, "Team Leader:  Siddhant", sz=16, color=LIGHT_GRAY)
tb(s, 7, 4.5, 5.8, 0.9, "Problem Statement: Track 03 — AI for Communities\nThe 'Last Mile' CSR Funding Gap (90% Misallocation)", sz=16, color=ACCENT_TEAL, align=PP_ALIGN.RIGHT)
footer(s)


# ═══════════════════════════════════════════════════
# SLIDE 2: BRIEF ABOUT THE IDEA
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); heading(s, "Brief About the Idea")

# Left: Problem
tb(s, 0.5, 1.2, 6, 0.5, "THE PROBLEM", sz=18, color=RED, bold=True)
bullets(s, 0.5, 1.8, 6, 3.2, [
    "India's CSR mandate (Section 135): Rs.38,000 Cr/year projected",
    "90% of these funds flow to the top 50 well-resourced NGOs",
    "3 Million+ civil society orgs are locked out of this capital",
    "Triple Barrier: Compliance paperwork + English fluency + Discovery",
    "A village schoolteacher feeding 200 children can't navigate this",
], sz=13)

# Right: Solution
tb(s, 7, 1.2, 6, 0.5, "THE SOLUTION", sz=18, color=GREEN, bold=True)
bullets(s, 7, 1.8, 5.8, 3.2, [
    "NidhiAI = AI-powered Fundraising Officer for small NGOs",
    "Upload docs > Compliance check > Grant matching > Proposal draft",
    "Uses Amazon Bedrock Multi-Agent Supervisor Pattern",
    "4 specialized agents orchestrated automatically as a team",
    "30-day manual process compressed into a 10-minute automated one",
], sz=13)

# Bottom stat bar
card(s, 0.5, 5.5, 12.333, 1.2, LIGHT_BG)
tb(s, 0.8, 5.65, 3.5, 0.9, "38,000 Cr\nCSR Market", sz=20, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
tb(s, 4.5, 5.65, 3.5, 0.9, "3 Million+\nRegistered NGOs", sz=20, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
tb(s, 8.5, 5.65, 4, 0.9, "90% Funds\nMisallocated", sz=20, color=RED, bold=True, align=PP_ALIGN.CENTER)
footer(s)


# ═══════════════════════════════════════════════════
# SLIDE 3: WHY AI IS REQUIRED
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); heading(s, "Why AI is Required + AWS Services in Architecture")

qa = [
    ("Why AI\nis Required", "This isn't a CRUD app.\n\n12A/80G certs are scanned\nPDFs with no API. Only\nTextract + an LLM agent\ncan extract dates reliably.\n\nKeyword search can't map\n'village water project' to\n'Schedule VII Rural Infra'.\nSemantic embeddings are\nstrictly required.\n\nNGOs can't afford Rs.50K\ngrant writers. Claude does\nit in 18 seconds."),
    ("How AWS Services\nAre Used", "Bedrock Agents: 5 agents\n(1 Supervisor + 4 sub-agents)\n\nBedrock KBs: 3 knowledge\nbases on OpenSearch Serverless\n\nTitan Embeddings V2 for\nvector search across KBs\n\nTextract: Document OCR\non legal certificates\n\nLambda x4: Serverless\ncompute for each action\n\nS3 x5 + DynamoDB x4"),
    ("What Value\nAI Adds", "Without AI:\n  30 days, Rs.50K per proposal\n  Manual search across websites\n  CA fee for compliance check\n\nWith NidhiAI:\n  10 minutes, Rs.10 per proposal\n  AI-matched across 5+ corps\n  Instant automated validation\n\n5000x cost reduction.\n200x faster delivery."),
]

for i, (title, desc) in enumerate(qa):
    x = 0.5 + i * 4.2
    card(s, x, 1.3, 3.8, 5.2, LIGHT_BG, ACCENT_TEAL)
    tb(s, x + 0.2, 1.45, 3.4, 0.8, title, sz=15, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, x + 0.3, 2.4, 3.2, 3.8, desc, sz=11, color=LIGHT_GRAY)

footer(s)


# ═══════════════════════════════════════════════════
# SLIDE 4: LIST OF FEATURES
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); heading(s, "Key Features of the Prototype")

features = [
    ("Compliance Scanner", "Upload 12A, 80G, CSR-1 certs.\nTextract scans, Bedrock validates.\nInstant red/green status."),
    ("AI Grant Discovery", "Semantic search across 5 corporate\nCSR programs. Ranked matches\nwith relevance scores (93%, 88%)."),
    ("Proposal Generator", "One-click formal 5-page proposals\ntailored to each grant's criteria.\nPDF output ready to send."),
    ("Impact Reports", "Auto-generates quarterly reports\nwith AI-computed metrics from\nplatform activity data."),
    ("CSR Assistant Chatbot", "Knowledge Base of Indian CSR laws\n(Section 135, Schedule VII, FCRA).\nAnswers regulatory questions."),
    ("Agent Trace Panel", "See the Supervisor's thinking live.\nShows which sub-agent was called,\nwhat data was retrieved."),
]

colors = [BLUE, ACCENT_ORANGE, PURPLE, GREEN, ACCENT_TEAL, RGBColor(0xFF, 0x66, 0x99)]
for i, ((title, desc), color) in enumerate(zip(features, colors)):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 4.3
    y = 1.3 + row * 3.0
    card(s, x, y, 3.9, 2.6, LIGHT_BG)
    # Color accent bar
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(2.6))
    accent.fill.solid(); accent.fill.fore_color.rgb = color; accent.line.fill.background()
    tb(s, x + 0.3, y + 0.2, 3.4, 0.4, title, sz=15, color=color, bold=True)
    tb(s, x + 0.3, y + 0.7, 3.4, 1.6, desc, sz=11, color=LIGHT_GRAY)

footer(s)


# ═══════════════════════════════════════════════════
# SLIDE 5: PROCESS FLOW
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); heading(s, "Process Flow: End-to-End User Journey")

steps = [
    ("👤 NGO Leader\nUploads Docs", ACCENT_TEAL),
    ("⚖️ Compliance\nAgent Scans", BLUE),
    ("🔍 Scout Agent\nFinds Grants", ACCENT_ORANGE),
    ("📝 Proposal Agent\nWrites PDF", PURPLE),
    ("📄 NGO Gets\nReady Proposal", GREEN),
]

sw = 2.0; gap = 0.6
total = len(steps) * sw + (len(steps)-1) * gap
sx = (13.333 - total) / 2
for i, (label, color) in enumerate(steps):
    x = sx + i * (sw + gap)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(sw), Inches(1.8))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(13); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    if i < len(steps)-1:
        ax = x + sw
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(ax + 0.05), Inches(2.65), Inches(gap - 0.1), Inches(0.5))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = ACCENT_TEAL; arrow.line.fill.background()

card(s, 3, 4.3, 7.333, 0.7, LIGHT_BG, border=ACCENT_ORANGE, text="🤖 Orchestrated by SUPERVISOR AGENT (Bedrock Supervisor Pattern)", sz=14, fc=ACCENT_ORANGE)

# Agent IDs row
agent_info = [
    ("Compliance Agent", "GFOH0VN84S"),
    ("Grant Scout Agent", "KRPTCZMV4Z"),
    ("Proposal Agent", "23TIS55Z0P"),
    ("Impact Agent", "PHLYZYB11A"),
    ("Supervisor Agent", "HB82HPMIA3"),
]
for i, (name, aid) in enumerate(agent_info):
    x = 0.5 + i * 2.55
    tb(s, x, 5.3, 2.3, 0.3, name, sz=10, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, x, 5.6, 2.3, 0.3, f"ID: {aid}", sz=9, color=DIM, align=PP_ALIGN.CENTER)

tb(s, 0.5, 6.2, 12, 0.4, "All 5 agents status: PREPARED  |  Region: ap-south-1 (Mumbai)  |  100% Serverless", sz=12, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
footer(s)


# ═══════════════════════════════════════════════════
# SLIDE 6: WIREFRAMES / SCREENSHOTS
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); heading(s, "Wireframes / Screenshots of the Working Prototype")

add_img(s, 'dashboard_1772968871753.png', 0.3, 1.2, 6.3, 3.8)
add_img(s, 'grants_1772969005674.png', 6.8, 1.2, 6.3, 3.8)

tb(s, 0.3, 5.1, 6.3, 0.4, "Dashboard — Live stats, compliance status, quick actions", sz=11, color=DIM, align=PP_ALIGN.CENTER)
tb(s, 6.8, 5.1, 6.3, 0.4, "Grant Discovery — AI-matched CSR opportunities with scores", sz=11, color=DIM, align=PP_ALIGN.CENTER)

# Bottom row
add_img(s, 'compliance_1772968882563.png', 0.3, 5.6, 4.1, 1.7)
add_img(s, 'proposals_1772969030004.png', 4.6, 5.6, 4.1, 1.7)
footer(s, "All screenshots from the LIVE running prototype on localhost:3000")


# ═══════════════════════════════════════════════════
# SLIDE 7: ARCHITECTURE DIAGRAM
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); heading(s, "Architecture Diagram: AWS Services")

# Top row
card(s, 0.3, 1.3, 2.5, 0.9, ACCENT_TEAL, text="NGO User\n(Next.js Dashboard)", sz=11)
card(s, 3.2, 1.3, 2.5, 0.9, RGBColor(0x44, 0x66, 0xBB), text="API Gateway\n+ Cognito Auth", sz=10)
card(s, 6.1, 1.3, 2.5, 0.9, RGBColor(0x22, 0x55, 0x99), text="Amazon S3\n(5 Buckets)", sz=10)

# Supervisor
card(s, 3.2, 2.6, 5.5, 1.0, ACCENT_ORANGE, text="BEDROCK SUPERVISOR AGENT (HB82HPMIA3)", sz=14)

# 4 Sub-agents
agents_data = [
    ("Compliance\n(GFOH0VN)", BLUE),
    ("Grant Scout\n(KRPTCZMV)", RGBColor(0xFF, 0x88, 0x00)),
    ("Proposal\n(23TIS55Z)", PURPLE),
    ("Impact\n(PHLYZYB1)", GREEN),
]
for i, (name, color) in enumerate(agents_data):
    x = 0.3 + i * 3.2
    card(s, x, 4.1, 2.8, 0.9, color, text=name, sz=10)

# Right: Service stack
services = [
    "Amazon Textract",
    "KB: CSR Laws (WUIYUPN7I2)",
    "KB: CSR Opps (UMGDGC60BJ)",
    "KB: Proposals (GXLH46FI6K)",
]
for i, svc in enumerate(services):
    card(s, 9.5, 1.3 + i * 0.85, 3.5, 0.7, RGBColor(0x22, 0x55, 0x99), text=svc, sz=9)

# Lambda row
lambdas = ["scan-documents", "match-grants", "generate-pdf", "generate-report"]
for i, fn in enumerate(lambdas):
    x = 0.3 + i * 3.2
    card(s, x, 5.4, 2.8, 0.55, RGBColor(0x22, 0x55, 0x99), text=f"Lambda: nidhiai-{fn}", sz=8)

# Bottom
card(s, 3.2, 6.1, 3.5, 0.55, RGBColor(0x22, 0x55, 0x99), text="Amazon DynamoDB (4 tables)", sz=9)
card(s, 7, 6.1, 3.5, 0.55, RGBColor(0x22, 0x55, 0x99), text="OpenSearch Serverless (Vector DB)", sz=9)

tb(s, 0.5, 6.8, 12, 0.4, "All services in ap-south-1 (Mumbai)  |  100% Serverless  |  12 AWS Services", sz=11, color=DIM, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 8: TECHNOLOGIES
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); heading(s, "Technologies Utilized in the Solution")

cats = [
    ("AI & ML Layer", [
        "Amazon Bedrock Agents (Supervisor Pattern, 5 agents)",
        "Amazon Bedrock Knowledge Bases (3 KBs, ACTIVE)",
        "Amazon Titan Text Embeddings V2",
        "Foundation Models: Claude 3.5 Sonnet / Amazon Nova 2 Lite",
    ]),
    ("Document Intelligence", [
        "Amazon Textract (OCR + Form Extraction)",
        "Scans 12A, 80G, CSR-1 govt certificates",
        "Extracts registration numbers and expiry dates",
        "Available natively in ap-south-1 (Mumbai)",
    ]),
    ("Compute & Storage", [
        "AWS Lambda x4 (Python 3.12, 256MB, 90s timeout)",
        "Amazon S3 x5 buckets (docs, PDFs, 3 KB data sources)",
        "Amazon OpenSearch Serverless (vector database)",
        "Amazon API Gateway (REST) + Amazon DynamoDB x4",
    ]),
    ("Frontend & Auth", [
        "Next.js 16 (React 19, Server Components)",
        "TailwindCSS + custom dark/light theme system",
        "Amazon Cognito (user authentication & sessions)",
        "Deployed on Vercel (CDN, Edge Functions)",
    ]),
]

for i, (cat, items) in enumerate(cats):
    col = i % 2; row = i // 2
    x = 0.5 + col * 6.5; y = 1.3 + row * 3.0
    card(s, x, y, 6.0, 2.7, LIGHT_BG, ACCENT_TEAL)
    tb(s, x + 0.2, y + 0.15, 5.5, 0.4, cat, sz=15, color=ACCENT_ORANGE, bold=True)
    bullets(s, x + 0.2, y + 0.6, 5.5, 1.8, items, sz=11)

footer(s)


# ═══════════════════════════════════════════════════
# SLIDE 9: COST
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); heading(s, "Estimated Implementation Cost")

# Left: Unit Economics
card(s, 0.5, 1.3, 6.0, 5.3, LIGHT_BG, ACCENT_TEAL)
tb(s, 0.7, 1.4, 5.5, 0.4, "Unit Economics: NidhiAI vs Traditional", sz=15, color=ACCENT_ORANGE, bold=True)

comparisons = [
    ("Metric", "Traditional", "NidhiAI", ACCENT_ORANGE, ACCENT_ORANGE, ACCENT_ORANGE),
    ("Cost per proposal", "Rs.50,000+", "< Rs.10", WHITE, RED, GREEN),
    ("Time per proposal", "2-4 weeks", "< 10 minutes", WHITE, RED, GREEN),
    ("Grant discovery", "Manual search", "AI-matched (5+ corps)", WHITE, RED, GREEN),
    ("Compliance check", "CA fees Rs.5,000+", "Automated via Textract", WHITE, RED, GREEN),
    ("Scaling to 100 NGOs", "50 grant writers", "Same AI system", WHITE, RED, GREEN),
]
for i, (m, t, n, mc, tc, nc) in enumerate(comparisons):
    y = 2.0 + i * 0.52
    bg = RGBColor(0x1A, 0x3A, 0x5A) if i == 0 else (DARK_BG if i % 2 == 0 else LIGHT_BG)
    row = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y), Inches(5.5), Inches(0.47))
    row.fill.solid(); row.fill.fore_color.rgb = bg; row.line.fill.background()
    tb(s, 0.8, y + 0.04, 1.8, 0.35, m, sz=11, color=mc, bold=(i==0))
    tb(s, 2.7, y + 0.04, 1.5, 0.35, t, sz=11, color=tc, bold=(i==0), align=PP_ALIGN.CENTER)
    tb(s, 4.3, y + 0.04, 1.7, 0.35, n, sz=11, color=nc, bold=(i==0), align=PP_ALIGN.CENTER)

# Savings callout
savings = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.2), Inches(5.5), Inches(0.9))
savings.fill.solid(); savings.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x1A)
savings.line.color.rgb = GREEN; savings.line.width = Pt(1.5)
tb(s, 0.9, 5.3, 5.1, 0.7, "5000x cheaper per proposal  |  200x faster delivery\nPrototype spend to date: < Rs.8,000", sz=12, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

# Right: Revenue Model
card(s, 7.0, 1.3, 5.8, 5.3, LIGHT_BG, ACCENT_ORANGE)
tb(s, 7.2, 1.4, 5.3, 0.4, "Revenue Model (100 NGOs)", sz=15, color=ACCENT_ORANGE, bold=True)

rev = [
    ("Free Tier", "5 proposals/month (Rs.0)"),
    ("Premium Plan", "Unlimited (Rs.999/NGO/mo)"),
    ("100 NGOs on Premium", "Rs.99,900/month revenue"),
    ("AWS Infra Cost", "~Rs.17,000/month (~$204)"),
    ("Gross Margin", "~83%"),
    ("Break-even", "18 paying NGOs"),
    ("B2B: CSR Consultants", "Rs.9,999/month enterprise"),
    ("Target Market", "3M+ civil society orgs"),
]
for i, (m, v) in enumerate(rev):
    y = 2.0 + i * 0.5
    mc = ACCENT_ORANGE if i == 0 else WHITE
    vc = ACCENT_ORANGE if i == 0 else (GREEN if i == 4 else WHITE)
    table_row(s, 7.2, y, 5.3, m, v, i, mc, vc, i == 0)

footer(s)


# ═══════════════════════════════════════════════════
# SLIDE 10: MORE PROTOTYPE SNAPSHOTS
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); heading(s, "Snapshots of the Working Prototype")

add_img(s, 'chatbot_final_1772969216606.png', 0.3, 1.2, 6.3, 3.5)
add_img(s, 'landing_page_1772968848332.png', 6.8, 1.2, 6.3, 3.5)

tb(s, 0.3, 4.8, 6.3, 0.4, "CSR Assistant Chatbot — answers regulatory questions via KB", sz=11, color=DIM, align=PP_ALIGN.CENTER)
tb(s, 6.8, 4.8, 6.3, 0.4, "Landing Page — Initialize demo, sign in/up", sz=11, color=DIM, align=PP_ALIGN.CENTER)

add_img(s, 'upload_page_fixed_1772819311612.png', 0.3, 5.4, 4.1, 1.5)
add_img(s, 'compliance_1772968882563.png', 4.6, 5.4, 4.1, 1.5)
# Bottom labels
footer(s, "All screenshots from the LIVE running prototype  |  Dark theme, responsive, production-quality UI")


# ═══════════════════════════════════════════════════
# SLIDE 11: PERFORMANCE
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); heading(s, "Prototype Performance Report / Benchmarking")

# Left: Latency metrics
card(s, 0.5, 1.3, 6.0, 4.0, LIGHT_BG, ACCENT_TEAL)
tb(s, 0.7, 1.45, 5.5, 0.4, "Latency (Live on AWS ap-south-1)", sz=15, color=ACCENT_ORANGE, bold=True)

perf = [
    ("Document Upload to S3", "~1.2 seconds"),
    ("Compliance Scan (Textract + Agent)", "4-6 seconds"),
    ("Grant Matching (Titan V2 + AOSS)", "1.8 seconds"),
    ("Proposal Generation (RAG + Claude)", "15-20s (streamed)"),
    ("KB Chatbot Response", "2-4 seconds"),
    ("Lambda Cold Start", "< 800ms"),
]
for i, (m, v) in enumerate(perf):
    table_row(s, 0.7, 2.0 + i * 0.47, 5.5, m, v, i, WHITE, ACCENT_TEAL, i == 0)

# Right: Infrastructure status
card(s, 7.0, 1.3, 5.8, 4.0, LIGHT_BG, ACCENT_ORANGE)
tb(s, 7.2, 1.45, 5.3, 0.4, "Infrastructure Status", sz=15, color=ACCENT_ORANGE, bold=True)

infra = [
    ("Bedrock Agents (5/5)", "PREPARED", GREEN),
    ("Knowledge Bases (3/3)", "ACTIVE", GREEN),
    ("KB Sync (18 docs)", "0 failures", GREEN),
    ("Lambda Functions (4/4)", "Active", GREEN),
    ("System Availability", "100%", GREEN),
    ("Concurrent Scaling", "Auto (Lambda)", ACCENT_TEAL),
]
for i, (m, v, vc) in enumerate(infra):
    y = 2.0 + i * 0.47
    bg = RGBColor(0x1A, 0x3A, 0x5A) if i == 0 else (DARK_BG if i % 2 == 0 else LIGHT_BG)
    row = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(y), Inches(5.3), Inches(0.42))
    row.fill.solid(); row.fill.fore_color.rgb = bg; row.line.fill.background()
    tb(s, 7.3, y + 0.03, 2.8, 0.35, m, sz=11, color=WHITE)
    tb(s, 10.2, y + 0.03, 2.1, 0.35, v, sz=11, color=vc, bold=True, align=PP_ALIGN.RIGHT)

# Bottom fact bar
card(s, 0.5, 5.6, 12.333, 1.0, LIGHT_BG)
tb(s, 0.8, 5.75, 4, 0.7, "4 Lambda Functions\nPython 3.12, 256MB", sz=14, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
tb(s, 4.8, 5.75, 4, 0.7, "5 S3 Buckets\n4 DynamoDB Tables", sz=14, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
tb(s, 8.8, 5.75, 4, 0.7, "3 Knowledge Bases\n18 Documents Indexed", sz=14, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
footer(s)


# ═══════════════════════════════════════════════════
# SLIDE 12: FUTURE DEVELOPMENT
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); heading(s, "Additional Details / Future Development")

# Left: Near-term
card(s, 0.5, 1.3, 6.0, 5.3, LIGHT_BG, ACCENT_TEAL)
tb(s, 0.7, 1.45, 5.5, 0.4, "Near-Term (3-6 months)", sz=15, color=ACCENT_ORANGE, bold=True)
bullets(s, 0.7, 2.0, 5.5, 3.5, [
    "Multilingual voice input via Amazon Transcribe — NGO founders",
    "  narrate in Hindi, AI converts to English corporate proposals",
    "",
    "Expand Grant KB from 5 corporations to 50+ (Tata, Infosys,",
    "  HDFC, Reliance, Wipro, Mahindra, Adani, etc.)",
    "",
    "Direct submission to corporate CSR portals (API integration)",
    "",
    "Automated compliance reminders (12A/80G renewal alerts)",
    "",
    "Mobile-first PWA for low-bandwidth rural areas",
], sz=11)

# Right: Long-term vision
card(s, 7.0, 1.3, 5.8, 5.3, LIGHT_BG, ACCENT_ORANGE)
tb(s, 7.2, 1.45, 5.3, 0.4, "Long-Term Vision", sz=15, color=ACCENT_ORANGE, bold=True)
bullets(s, 7.2, 2.0, 5.3, 2.5, [
    "NidhiAI becomes the default operating system for",
    "  NGO administration in India",
    "",
    "Government scheme matching (PM-KISAN, MGNREGA)",
    "",
    "Amazon Q Business integration for internal CSR portals",
    "",
    "WhatsApp bot for zero-tech-literacy users",
], sz=11)

# Impact box
card(s, 7.2, 4.6, 5.3, 1.6, RGBColor(0x0A, 0x2A, 0x1A), GREEN)
tb(s, 7.4, 4.7, 5, 1.3,
    "If just 10% of Rs.38,000 Cr CSR pool reaches\nsmaller NGOs through NidhiAI:\n\nRs.3,800 Crores unlocked for communities\nthat need it most.",
    sz=12, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

footer(s)


# ═══════════════════════════════════════════════════
# SLIDE 13: PROTOTYPE ASSETS
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s); heading(s, "Prototype Assets")

card(s, 2, 1.8, 9.333, 1.5, LIGHT_BG, ACCENT_TEAL)
tb(s, 2.3, 2.0, 8.7, 0.5, "GitHub Public Repository", sz=18, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
tb(s, 2.3, 2.55, 8.7, 0.5, "https://github.com/Sid-V5/NidhiAi", sz=16, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

card(s, 2, 3.6, 9.333, 1.5, LIGHT_BG, ACCENT_ORANGE)
tb(s, 2.3, 3.8, 8.7, 0.5, "Demo Video (Max 3 Minutes)", sz=18, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
tb(s, 2.3, 4.35, 8.7, 0.5, "[Insert YouTube / Google Drive Link Here]", sz=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

card(s, 2, 5.4, 9.333, 1.2, LIGHT_BG, GREEN)
tb(s, 2.3, 5.55, 8.7, 0.5, "Working MVP Link", sz=18, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
tb(s, 2.3, 6.0, 8.7, 0.5, "[Insert Vercel Deployment URL Here]", sz=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

footer(s)


# ═══════════════════════════════════════════════════
# SLIDE 14: THANK YOU
# ═══════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
tb(s, 1, 1.5, 11, 1.0, "NidhiAI (निधिAI)", sz=48, color=ACCENT_ORANGE, bold=True, align=PP_ALIGN.CENTER)
tb(s, 1.5, 3.0, 10, 0.8, "Let's Fund the Changemakers of Bharat.", sz=28, color=WHITE, align=PP_ALIGN.CENTER)

div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.2), Inches(5.333), Inches(0.03))
div.fill.solid(); div.fill.fore_color.rgb = ACCENT_TEAL; div.line.fill.background()

tb(s, 1, 4.7, 11, 0.5, "Team: NidhiAI  |  Leader: Siddhant", sz=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
tb(s, 1, 5.3, 11, 0.5, "Track 03: AI for Communities, Access & Public Impact", sz=14, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
tb(s, 1, 5.9, 11, 0.5, "Powered by AWS Bedrock  •  Amazon Textract  •  Amazon OpenSearch", sz=13, color=DIM, align=PP_ALIGN.CENTER)
footer(s, "AI for Bharat Hackathon 2026  |  Thank You")


# ═══════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════
out = r'c:\Users\Siddhant\NidhiAi\NidhiAI_Prototype_Submission.pptx'
prs.save(out)
print(f"Saved to: {out}")
print(f"Total slides: {len(prs.slides)}")
